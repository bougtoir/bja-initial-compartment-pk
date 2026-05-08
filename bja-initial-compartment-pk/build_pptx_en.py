#!/usr/bin/env python3
"""Generate English PPTX figure set for RAPM submission.

One slide per figure, with title and caption.
Widescreen (13.333 x 7.5 inches).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Widescreen dimensions
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

fig_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'figures')

# Figure data
figures = [
    {
        'file': 'figure1_compartment_models.png',
        'title': 'Figure 1. Comparison of Pharmacokinetic Model Structures',
        'caption': (
            'Comparison of pharmacokinetic model structures for intravenous and '
            'regional administration. (A) Standard three-compartment intravenous model '
            'with drug entering the central plasma compartment (V1) directly. '
            '(B) Depot-augmented model for regional anesthesia with drug deposited into '
            'a tissue depot compartment, from which first-order absorption (rate constant '
            'ka) governs entry into the central compartment. V1, central compartment '
            '(plasma); V2, rapidly equilibrating peripheral compartment (vessel-rich '
            'tissues); V3, slowly equilibrating peripheral compartment (vessel-poor '
            'tissues); CL, clearance; ka, absorption rate constant; F, bioavailability.'
        ),
    },
    {
        'file': 'figure2_pk_simulation.png',
        'title': 'Figure 2. Simulated Plasma Concentration\u2013Time Profiles',
        'caption': (
            'Simulated plasma concentration\u2013time profiles demonstrating the effect of '
            'route of administration on local anesthetic pharmacokinetics. Blue solid '
            'line: intravenous bolus (standard three-compartment model). Red dashed line: '
            'rapid absorption (ka = 0.1 min\u207b\u00b9, representing highly vascular '
            'injection site). Green dash-dot line: slow absorption (ka = 0.03 min\u207b\u00b9, '
            'representing fascial plane deposition). Shaded regions indicate traditional '
            'monitoring window (blue, 0\u201330 min) and true Tmax window for slow absorption '
            '(green). Note the marked differences in Cmax and Tmax depending on the '
            'absorption rate.'
        ),
    },
    {
        'file': 'figure3_workflow.png',
        'title': 'Figure 3. Route-Specific PBPK Model Development Workflow',
        'caption': (
            'Proposed workflow for development and validation of route-specific '
            'pharmacokinetic models for regional anesthesia. Step 1: Published population '
            'pharmacokinetic data provide block-type-specific absorption parameters. '
            'Step 2: PBPK model configuration using PK-Sim/MoBi with route-specific '
            'initial conditions. Step 3: Model validation against independent clinical '
            'pharmacokinetic datasets. Step 4: If validated, potential integration into '
            'clinical decision support tools as a complement to existing dose limits.'
        ),
    },
]

blank_layout = prs.slide_layouts[6]  # Blank layout

for fig_info in figures:
    slide = prs.slides.add_slide(blank_layout)

    # Title at top
    title_left = Inches(0.5)
    title_top = Inches(0.3)
    title_width = Inches(12.3)
    title_height = Inches(0.7)
    txBox = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = fig_info['title']
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.LEFT

    # Image centered
    fig_path = os.path.join(fig_dir, fig_info['file'])
    if os.path.exists(fig_path):
        img_width = Inches(10.5)
        img_left = (SLIDE_WIDTH - img_width) / 2
        img_top = Inches(1.1)
        slide.shapes.add_picture(fig_path, int(img_left), int(img_top), width=int(img_width))

    # Caption at bottom
    cap_left = Inches(0.5)
    cap_top = Inches(5.8)
    cap_width = Inches(12.3)
    cap_height = Inches(1.5)
    txBox2 = slide.shapes.add_textbox(cap_left, cap_top, cap_width, cap_height)
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = fig_info['caption']
    p2.font.size = Pt(11)
    p2.font.name = 'Arial'
    p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p2.alignment = PP_ALIGN.LEFT

# Save
out_dir = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(out_dir, 'RAPM_Figures_English.pptx')
prs.save(out_path)
print(f'Saved: {out_path}')
