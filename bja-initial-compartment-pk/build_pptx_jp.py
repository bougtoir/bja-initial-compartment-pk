#!/usr/bin/env python3
"""Generate Japanese PPTX with 1 figure/table per slide.
- Figures 1, 2: embedded as PNG images (code-generated)
- Figure 3: editable PowerPoint shapes (workflow diagram)
- Figure 4: editable PowerPoint shapes (AIMS schematic)
- Table 1: editable PowerPoint table
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x21, 0x96, 0xF3)
RED = RGBColor(0xF4, 0x43, 0x36)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
ORANGE = RGBColor(0xFF, 0x98, 0x00)
GREY = RGBColor(0x9E, 0x9E, 0x9E)
PURPLE = RGBColor(0x9C, 0x27, 0xB0)
DARK = RGBColor(0x37, 0x47, 0x4F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
LIGHT_BLUE = RGBColor(0xE3, 0xF2, 0xFD)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
LIGHT_RED = RGBColor(0xFC, 0xE4, 0xEC)
LIGHT_ORANGE = RGBColor(0xFF, 0xF3, 0xE0)
LIGHT_PURPLE = RGBColor(0xF3, 0xE5, 0xF5)
LIGHT_GREY = RGBColor(0xF5, 0xF5, 0xF5)
TEAL = RGBColor(0x00, 0x69, 0x7C)

def add_title_textbox(slide, text, left, top, width, height, font_size=24, bold=True, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    return txBox

def add_textbox(slide, text, left, top, width, height, font_size=11, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None, text='', font_size=11, font_color=BLACK, bold=False, text_align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = text_align
    if text:
        lines = text.split('\n')
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(font_size)
            p.font.color.rgb = font_color
            p.font.bold = bold
            p.alignment = text_align
    shape.text_frame.auto_size = None
    shape.text_frame.word_wrap = True
    return shape

def add_arrow_shape(slide, left, top, width, height, color=DARK):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_down_arrow(slide, left, top, width, height, color=DARK):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ============================================================
# SLIDE 1: Figure 1 (PNG image)
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_textbox(slide1, '\u56f31. 3\u3064\u306e\u81e8\u5e8a\u30b7\u30ca\u30ea\u30aa\u306b\u304a\u3051\u308b\u30b3\u30f3\u30d1\u30fc\u30c8\u30e1\u30f3\u30c8\u30e2\u30c7\u30eb\u306e\u6bd4\u8f03',
                  Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6), font_size=20)
fig1_path = '/home/ubuntu/manuscript/figures/figure1_compartment_models.png'
if os.path.exists(fig1_path):
    slide1.shapes.add_picture(fig1_path, Inches(0.5), Inches(1.0), Inches(12.3), Inches(5.5))
add_textbox(slide1, '(A) \u5f93\u6765\u306eIV\u30e2\u30c7\u30eb  |  (B) \u6210\u529f\u3057\u305f\u533a\u57df\u30d6\u30ed\u30c3\u30af (BPT\u30b9\u30bf\u30fc\u30c8)  |  (C) \u4e0d\u6210\u529f\u306e\u30d6\u30ed\u30c3\u30af / \u8840\u7ba1\u5185\u6ce8\u5165 (\u8840\u6f3f/BRT\u30b9\u30bf\u30fc\u30c8)',
            Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5), font_size=12, color=DARK, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Figure 2 (PNG image)
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_textbox(slide2, '\u56f32. \u521d\u671f\u30b3\u30f3\u30d1\u30fc\u30c8\u30e1\u30f3\u30c8\u5225\u306e\u30b7\u30df\u30e5\u30ec\u30fc\u30b7\u30e7\u30f3\u8840\u6f3f\u6fc3\u5ea6\u2013\u6642\u9593\u30d7\u30ed\u30d5\u30a1\u30a4\u30eb',
                  Inches(0.5), Inches(0.1), Inches(12.3), Inches(0.8), font_size=20)
fig2_path = '/home/ubuntu/manuscript/figures/figure2_pk_simulation.png'
if os.path.exists(fig2_path):
    slide2.shapes.add_picture(fig2_path, Inches(1.5), Inches(1.0), Inches(10.3), Inches(6.0))

# ============================================================
# SLIDE 3: Figure 3 (EDITABLE workflow)
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_textbox(slide3, '\u56f33. \u533a\u57df\u9ebb\u9154\u306b\u304a\u3051\u308bPBPK\u30d9\u30fc\u30b9\u306e\u6975\u91cf\u518d\u8a55\u4fa1\u30ef\u30fc\u30af\u30d5\u30ed\u30fc',
                  Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.6), font_size=20)

step_w = Inches(3.2)
step_h = Inches(1.0)
y_top = Inches(1.0)

# Step 1
add_rounded_rect(slide3, Inches(0.5), y_top, step_w, step_h,
                 RGBColor(0x5C, 0x6B, 0xC0), line_color=RGBColor(0x3F, 0x51, 0xB5),
                 text='\u30b9\u30c6\u30c3\u30d71\n\u81e8\u5e8a\u8a55\u4fa1', font_size=14, font_color=WHITE, bold=True)

add_arrow_shape(slide3, Inches(3.8), Inches(1.25), Inches(0.7), Inches(0.3), color=GREY)

# Step 2
add_rounded_rect(slide3, Inches(4.6), y_top, step_w, step_h,
                 ORANGE, line_color=RGBColor(0xE6, 0x51, 0x00),
                 text='\u30b9\u30c6\u30c3\u30d72\n\u521d\u671f\u30b3\u30f3\u30d1\u30fc\u30c8\u30e1\u30f3\u30c8\u9078\u629e', font_size=14, font_color=WHITE, bold=True)

add_arrow_shape(slide3, Inches(7.9), Inches(1.25), Inches(0.7), Inches(0.3), color=GREY)

# Step 3
add_rounded_rect(slide3, Inches(8.7), y_top, step_w, step_h,
                 PURPLE, line_color=RGBColor(0x6A, 0x1B, 0x9A),
                 text='\u30b9\u30c6\u30c3\u30d73\nPBPK\u30b7\u30df\u30e5\u30ec\u30fc\u30b7\u30e7\u30f3\n(PK-Sim / MoBi)', font_size=13, font_color=WHITE, bold=True)

# Down arrows
y_mid = Inches(2.3)
add_down_arrow(slide3, Inches(2.0), y_mid, Inches(0.3), Inches(0.5), color=GREEN)
add_down_arrow(slide3, Inches(6.0), y_mid, Inches(0.3), Inches(0.5), color=ORANGE)
add_down_arrow(slide3, Inches(10.0), y_mid, Inches(0.3), Inches(0.5), color=RED)

# Three scenario boxes
sc_w = Inches(3.5)
sc_h = Inches(1.8)
y_sc = Inches(3.0)

add_rounded_rect(slide3, Inches(0.3), y_sc, sc_w, sc_h,
                 LIGHT_GREEN, line_color=GREEN,
                 text='\u6210\u529f\u3057\u305f\u30d6\u30ed\u30c3\u30af\nBPT\u30b9\u30bf\u30fc\u30c8\n\n\u4f4e\u3044Cmax\n\u3088\u308a\u9ad8\u7528\u91cf\u304c\u5b89\u5168\u306a\u53ef\u80fd\u6027', font_size=12, font_color=RGBColor(0x2E, 0x7D, 0x32))

add_rounded_rect(slide3, Inches(4.5), y_sc, sc_w, sc_h,
                 LIGHT_ORANGE, line_color=ORANGE,
                 text='\u90e8\u5206\u7684\u30d6\u30ed\u30c3\u30af\n\u6df7\u5408\u30b9\u30bf\u30fc\u30c8\n\n\u4e2d\u9593\u7684Cmax', font_size=12, font_color=RGBColor(0xE6, 0x51, 0x00))

add_rounded_rect(slide3, Inches(8.7), y_sc, sc_w, sc_h,
                 LIGHT_RED, line_color=RED,
                 text='\u4e0d\u6210\u529f\u306e\u30d6\u30ed\u30c3\u30af\n\u8840\u6f3f/BRT\u30b9\u30bf\u30fc\u30c8\n\n\u9ad8\u3044Cmax\n\u3088\u308a\u4f4e\u7528\u91cf\u304c\u5fc5\u8981\u306a\u53ef\u80fd\u6027', font_size=12, font_color=RGBColor(0xC6, 0x28, 0x28))

# Down arrows to Step 4
add_down_arrow(slide3, Inches(2.0), Inches(5.0), Inches(0.3), Inches(0.5), color=GREEN)
add_down_arrow(slide3, Inches(6.0), Inches(5.0), Inches(0.3), Inches(0.5), color=ORANGE)
add_down_arrow(slide3, Inches(10.0), Inches(5.0), Inches(0.3), Inches(0.5), color=RED)

# Step 4
add_rounded_rect(slide3, Inches(0.5), Inches(5.7), Inches(12.0), Inches(1.3),
                 RGBColor(0x00, 0x69, 0x7C), line_color=TEAL,
                 text='\u30b9\u30c6\u30c3\u30d74\uff1a\u72b6\u6cc1\u4f9d\u5b58\u6027\u6975\u91cf\u63a8\u5968\n\n\u4e88\u6e2c\u8840\u6f3f\u6fc3\u5ea6\u30d7\u30ed\u30d5\u30a1\u30a4\u30eb\u306b\u57fa\u3065\u304f\u500b\u5225\u5316\u3055\u308c\u305f\u30b7\u30ca\u30ea\u30aa\u4f9d\u5b58\u6027\u6975\u91cf',
                 font_size=14, font_color=WHITE, bold=True)

add_textbox(slide3, '\u30d6\u30ed\u30c3\u30af\u7a2e\u5225\u3001\u90e8\u4f4d\u3001\u6210\u529f\u78ba\u7387\u3001\u60a3\u8005\u56e0\u5b50',
            Inches(0.3), Inches(2.1), Inches(3.0), Inches(0.4), font_size=9, color=GREY)
add_textbox(slide3, '\u30b7\u30df\u30e5\u30ec\u30fc\u30b7\u30e7\u30f3\u7d50\u679c',
            Inches(11.5), Inches(4.5), Inches(1.5), Inches(0.4), font_size=9, color=PURPLE)

# ============================================================
# SLIDE 4: Figure 4 (EDITABLE AIMS schematic)
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_textbox(slide4, '\u56f34. AIMS\u306b\u304a\u3051\u308b\u6295\u4e0e\u7d4c\u8def\u9069\u5fdc\u578bPKPD\u30b7\u30df\u30e5\u30ec\u30fc\u30b7\u30e7\u30f3',
                  Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.5), font_size=20)

# Outer AIMS box
aims_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(0.7), Inches(12.7), Inches(6.5))
aims_box.fill.solid()
aims_box.fill.fore_color.rgb = LIGHT_GREY
aims_box.line.color.rgb = DARK
aims_box.line.width = Pt(2)
tf = aims_box.text_frame
p = tf.paragraphs[0]
p.text = '\u9ebb\u9154\u60c5\u5831\u7ba1\u7406\u30b7\u30b9\u30c6\u30e0\uff08AIMS\uff09'
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DARK
p.alignment = PP_ALIGN.CENTER

# Drug Administration Record
add_rounded_rect(slide4, Inches(0.6), Inches(1.3), Inches(4.0), Inches(1.5),
                 LIGHT_BLUE, line_color=RGBColor(0x15, 0x65, 0xC0),
                 text='\u85ac\u5264\u6295\u4e0e\u8a18\u9332\n\n\u6295\u4e0e\u7d4c\u8def\uff1aIV / \u533a\u57df\u9ebb\u9154\uff08\u30d6\u30ed\u30c3\u30af\u7a2e\u5225\uff09\n\u85ac\u5264\uff1a\u30d6\u30d4\u30d0\u30ab\u30a4\u30f3 150 mg',
                 font_size=11, font_color=RGBColor(0x15, 0x65, 0xC0))

add_arrow_shape(slide4, Inches(4.7), Inches(1.8), Inches(0.8), Inches(0.3), color=DARK)

# Route Detection
add_rounded_rect(slide4, Inches(5.6), Inches(1.3), Inches(3.5), Inches(1.5),
                 LIGHT_ORANGE, line_color=RGBColor(0xE6, 0x51, 0x00),
                 text='\u6295\u4e0e\u7d4c\u8def\u691c\u51fa\n\nIF \u7d4c\u8def = IV:\n  \u2192 \u6a19\u6e96 3\u30b3\u30f3\u30d1\u30fc\u30c8\u30e1\u30f3\u30c8\nIF \u7d4c\u8def = \u533a\u57df\u9ebb\u9154:\n  \u2192 \u30c7\u30dd\u589e\u5f37\u30e2\u30c7\u30eb',
                 font_size=10, font_color=RGBColor(0xE6, 0x51, 0x00))

# IV Model
add_rounded_rect(slide4, Inches(0.6), Inches(3.2), Inches(4.0), Inches(1.8),
                 LIGHT_GREEN, line_color=RGBColor(0x2E, 0x7D, 0x32),
                 text='IV\u30e2\u30c7\u30eb\n\n3\u30b3\u30f3\u30d1\u30fc\u30c8\u30e1\u30f3\u30c8\uff08\u8840\u6f3f\u30b9\u30bf\u30fc\u30c8\uff09\nV1 \u2192 V2 (BRT) + V3 (BPT)\n\n\u6a19\u6e96 Marsh / Schnider / Eleveld',
                 font_size=11, font_color=RGBColor(0x2E, 0x7D, 0x32))

# Regional Model
add_rounded_rect(slide4, Inches(5.2), Inches(3.2), Inches(4.5), Inches(1.8),
                 LIGHT_RED, line_color=RGBColor(0xC6, 0x28, 0x28),
                 text='\u533a\u57df\u9ebb\u9154\u30e2\u30c7\u30eb\uff08\u6295\u4e0e\u7d4c\u8def\u9069\u5fdc\u578b\uff09\n\n\u30c7\u30dd + 3\u30b3\u30f3\u30d1\u30fc\u30c8\u30e1\u30f3\u30c8 (BPT\u30b9\u30bf\u30fc\u30c8)\nDepot(ka, F) \u2192 V1 \u2192 V2 + V3\n\n\u30d6\u30ed\u30c3\u30af\u7a2e\u5225\u7279\u7570\u7684ka\u5024\n(TAP, ESP, FNB, \u786c\u819c\u5916 \u7b49)',
                 font_size=10, font_color=RGBColor(0xC6, 0x28, 0x28))

# Block Success Feedback
add_rounded_rect(slide4, Inches(10.2), Inches(3.4), Inches(2.5), Inches(1.3),
                 LIGHT_PURPLE, line_color=RGBColor(0x6A, 0x1B, 0x9A),
                 text='\u30d6\u30ed\u30c3\u30af\u6210\u529f\n\u30d5\u30a3\u30fc\u30c9\u30d0\u30c3\u30af\n\n\u30ea\u30a2\u30eb\u30bf\u30a4\u30e0\u3067ka\u8abf\u6574',
                 font_size=10, font_color=RGBColor(0x6A, 0x1B, 0x9A))

# Arrows
add_down_arrow(slide4, Inches(2.4), Inches(2.85), Inches(0.25), Inches(0.35), color=RGBColor(0x2E, 0x7D, 0x32))
add_textbox(slide4, 'IV', Inches(2.0), Inches(2.85), Inches(0.5), Inches(0.3), font_size=10, bold=True, color=RGBColor(0x2E, 0x7D, 0x32))

add_down_arrow(slide4, Inches(7.3), Inches(2.85), Inches(0.25), Inches(0.35), color=RGBColor(0xC6, 0x28, 0x28))
add_textbox(slide4, '\u533a\u57df\u9ebb\u9154', Inches(7.6), Inches(2.85), Inches(1.2), Inches(0.3), font_size=10, bold=True, color=RGBColor(0xC6, 0x28, 0x28))

add_arrow_shape(slide4, Inches(9.8), Inches(3.9), Inches(0.4), Inches(0.2), color=RGBColor(0x6A, 0x1B, 0x9A))

# AIMS Display
add_rounded_rect(slide4, Inches(0.6), Inches(5.3), Inches(12.0), Inches(1.7),
                 WHITE, line_color=DARK, text='', font_size=12, font_color=DARK, bold=True)

add_textbox(slide4, '\u30ea\u30a2\u30eb\u30bf\u30a4\u30e0 AIMS \u30c7\u30a3\u30b9\u30d7\u30ec\u30a4', Inches(4.5), Inches(5.35), Inches(4.0), Inches(0.35),
            font_size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)

add_textbox(slide4, '\u4e88\u6e2c\u8840\u6f3f\u6fc3\u5ea6\n\n'
            '--- IV\u30e2\u30c7\u30eb\uff1a\u6025\u901f\u30d4\u30fc\u30af\u3001\u6025\u901f\u4f4e\u4e0b\n'
            '--- \u533a\u57df\u9ebb\u9154\u30e2\u30c7\u30eb\uff1a\u4f4e\u304f\u9045\u5ef6\u3057\u305fCmax\n'
            '--- \u6bd2\u6027\u95be\u5024\u30e9\u30a4\u30f3',
            Inches(0.8), Inches(5.7), Inches(4.5), Inches(1.2), font_size=10, color=DARK)

add_textbox(slide4, '\u7528\u91cf\u30ac\u30a4\u30c0\u30f3\u30b9\u30d1\u30cd\u30eb\n\n'
            '\u73fe\u5728\u7528\u91cf\uff1a\u30d6\u30d4\u30d0\u30ab\u30a4\u30f3 150 mg\n'
            '\u6295\u4e0e\u7d4c\u8def\uff1aTAP\u30d6\u30ed\u30c3\u30af\uff08\u4e21\u5074\uff09\n'
            '\u4e88\u6e2cCmax\uff1a0.8 \u00b5g/mL\uff08\u5b89\u5168\uff09\n'
            '\u6b8b\u4f59\u30de\u30fc\u30b8\u30f3\uff1a\u95be\u5024\u4e0b68%',
            Inches(6.5), Inches(5.7), Inches(5.5), Inches(1.2), font_size=10, color=DARK)

add_down_arrow(slide4, Inches(2.4), Inches(5.05), Inches(0.25), Inches(0.25), color=DARK)
add_down_arrow(slide4, Inches(7.3), Inches(5.05), Inches(0.25), Inches(0.25), color=DARK)

# ============================================================
# SLIDE 5: Table 1 (EDITABLE)
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_textbox(slide5, '\u88681. \u533a\u57df\u9ebb\u9154\u306b\u304a\u3051\u308b\u5c40\u6240\u9ebb\u9154\u85ac\u306e\u72b6\u6cc1\u4f9d\u5b58\u6027\u6975\u91cf\u306e\u679a\u7d44\u307f',
                  Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.8), font_size=20)

tbl_w = Inches(11.5)
tbl_h = Inches(4.5)
left = Inches(0.9)
top = Inches(1.3)

table_shape = slide5.shapes.add_table(6, 4, left, top, tbl_w, tbl_h)
table = table_shape.table

col_widths = [Inches(3.2), Inches(3.0), Inches(2.3), Inches(3.0)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w

headers = ['\u30b7\u30ca\u30ea\u30aa', '\u521d\u671f\u30b3\u30f3\u30d1\u30fc\u30c8\u30e1\u30f3\u30c8', '\u4e88\u6e2cCmax', '\u7528\u91cf\u8abf\u6574']

data = [
    ['\u6210\u529f\u3057\u305f\u30d6\u30ed\u30c3\u30af\n\uff08\u795e\u7d4c\u5468\u56f2/\u7b4b\u819c\u9762\uff09', 'BPT\uff08V3\uff09\n\u8840\u7ba1\u4e4f\u3057\u3044\u7d44\u7e54', '\u4f4e\u5024\u3001\u9045\u5ef6', '\u3088\u308a\u9ad8\u7528\u91cf\u304c\n\u5b89\u5168\u306a\u53ef\u80fd\u6027'],
    ['\u90e8\u5206\u7684\u30d6\u30ed\u30c3\u30af', '\u6df7\u5408\n\uff08BPT + BRT/\u8840\u6f3f\uff09', '\u4e2d\u9593', '\u6a19\u6e96\u7684\u6975\u91cf\n\u304c\u9069\u7528'],
    ['\u4e0d\u6210\u529f\u306e\u30d6\u30ed\u30c3\u30af\n\uff08\u7d44\u7e54\u30df\u30b9\u30d7\u30ec\u30a4\u30b9\u30e1\u30f3\u30c8\uff09', 'BRT\uff08V2\uff09\n\u8840\u7ba1\u8c4a\u5bcc\u7d44\u7e54', '\u4e2d\u7b49\u5ea6\uff5e\u9ad8\u5024\u3001\n\u65e9\u671f', '\u3088\u308a\u4f4e\u7528\u91cf\u304c\n\u5fc5\u8981\u306a\u53ef\u80fd\u6027'],
    ['\u8840\u7ba1\u5185\u6ce8\u5165', '\u8840\u6f3f\uff08V1\uff09', '\u975e\u5e38\u306b\u9ad8\u5024\u3001\n\u5373\u6642', '\u5f93\u6765\u306eIV\u6975\u91cf\n\u304c\u9069\u7528'],
]

row_colors = [LIGHT_GREEN, LIGHT_ORANGE, LIGHT_RED, RGBColor(0xFF, 0xCD, 0xD2), LIGHT_PURPLE]

for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.bold = True
        paragraph.font.color.rgb = WHITE
        paragraph.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

for r, row_data in enumerate(data):
    for c, val in enumerate(row_data):
        cell = table.cell(r + 1, c)
        cell.text = val
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = BLACK
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = row_colors[r]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

# Footnote about spinal exclusion
add_textbox(slide5, '注：脊髄くも膜下麻酔（脊椎麻酔）はCSF薬物動態の独自性、少量投与、および主に単回\n投与である手技の性質から、本枠組みから除外した。',
            Inches(0.9), Inches(6.0), Inches(11.5), Inches(0.7), font_size=10, color=GREY, align=PP_ALIGN.LEFT)

# Save
out = '/home/ubuntu/manuscript/BJA_Figures_Japanese.pptx'
prs.save(out)
print(f'Japanese PPTX saved: {out}')
